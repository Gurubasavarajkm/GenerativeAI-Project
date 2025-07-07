import pypdf
import docx
import openpyxl
import pptx
import os
from io import BytesIO

def extract_text_from_file(file_content: BytesIO, filename: str) -> str:
    _, file_extension = os.path.splitext(filename)
    text = ""
    
    if file_extension == '.pdf':
        reader = pypdf.PdfReader(file_content)
        for page in reader.pages:
            text += page.extract_text() or ""
    elif file_extension == '.docx':
        doc = docx.Document(file_content)
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif file_extension == '.xlsx':
        workbook = openpyxl.load_workbook(file_content)
        for sheet in workbook.sheetnames:
            worksheet = workbook[sheet]
            for row in worksheet.iter_rows():
                for cell in row:
                    if cell.value:
                        text += str(cell.value) + " "
                text += "\n"
    elif file_extension == '.pptx':
        presentation = pptx.Presentation(file_content)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    else:
        raise ValueError("Unsupported file type")
        
    return text